
import sys
from pptx import Presentation
from pptx.util import Inches

def extract_text(path):
    try:
        prs = Presentation(path)
        full_text = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            # Sort shapes by top position to maintain order
            shapes = sorted(slide.shapes, key=lambda shape: shape.top if hasattr(shape, 'top') else 0)
            
            for shape in shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                full_text.append(f"--- SLIDE {slide_num + 1} ---")
                full_text.append("\n".join(slide_text))
        
        return "\n".join(full_text)
    except Exception as e:
        return f"Error: {str(e)}"

if __name__ == "__main__":
    if len(sys.argv) > 1:
        print(extract_text(sys.argv[1]))
    else:
        print("No file path provided")
